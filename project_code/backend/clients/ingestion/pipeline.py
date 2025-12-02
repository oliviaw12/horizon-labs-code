"""Slide/document ingestion pipeline for vector search: extracts text from PPTX/PDF, chunks,
embeds via Google GenAI, and upserts to Pinecone with dimension validation."""

from __future__ import annotations

import io
import logging
import asyncio
from dataclasses import dataclass
from typing import Any, Dict, Iterator, List, Optional, Sequence

try:  # pragma: no cover - dependency optional during local testing
    from pypdf import PdfReader as _PdfReader  # type: ignore
except ModuleNotFoundError:  # pragma: no cover - executed when package missing
    _PdfReader = None  # type: ignore[assignment]

from clients.database.pinecone import PineconeRepository
from clients.llm.settings import Settings

logger = logging.getLogger(__name__)

PdfReader = _PdfReader


@dataclass
class SlideChunk:
    """A semantically meaningful chunk extracted from a document page or slide."""

    slide_number: int
    text: str
    slide_title: Optional[str]
    chunk_index: int
    source_type: str = "slide"

    def metadata(self) -> Dict[str, Any]:
        """Return metadata describing the origin of this chunk (slide/page, index, title)."""
        payload: Dict[str, Any] = {
            "slide_number": self.slide_number,
            "chunk_index": self.chunk_index,
            "slide_title": self.slide_title,
            "source_type": self.source_type,
        }
        if self.source_type == "page":
            payload["page_number"] = self.slide_number
        return payload


@dataclass
class IngestionResult:
    """Summary returned once the pipeline finishes ingesting a document."""

    document_id: str
    slide_count: int
    chunk_count: int
    namespace: str


class SlideExtractor:
    """Pulls raw text (and titles) from a slide deck."""

    def extract(self, file_bytes: bytes) -> List[SlideChunk]:
        try:
            from pptx import Presentation  # type: ignore
        except ModuleNotFoundError as exc:  # pragma: no cover - import guard
            raise RuntimeError(
                "python-pptx is required to ingest PowerPoint files. Install the dependency to continue."
            ) from exc

        presentation = Presentation(io.BytesIO(file_bytes))
        chunks: List[SlideChunk] = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            title = None
            if slide.shapes.title:
                title = slide.shapes.title.text.strip() or None

            text_fragments: List[str] = []
            for shape in slide.shapes:
                if not hasattr(shape, "text"):
                    continue
                payload = (shape.text or "").strip()
                if payload:
                    text_fragments.append(payload)

            slide_text = "\n".join(text_fragments).strip()
            if not slide_text:
                continue

            chunks.append(
                SlideChunk(
                    slide_number=slide_number,
                    text=slide_text,
                    slide_title=title,
                    chunk_index=0,
                    source_type="slide",
                )
            )
        return chunks


class PDFExtractor:
    """Extracts page-level text from a PDF document."""

    def extract(self, file_bytes: bytes) -> List[SlideChunk]:
        if PdfReader is None:
            raise RuntimeError(
                "pypdf is required to ingest PDF files. Install the dependency to continue."
            )

        reader = PdfReader(io.BytesIO(file_bytes))
        chunks: List[SlideChunk] = []
        for page_number, page in enumerate(reader.pages, start=1):
            try:
                page_text = (page.extract_text() or "").strip()
            except Exception:  # pragma: no cover - defensive guard for uncommon PDFs
                logger.exception("Failed extracting text from PDF page %s", page_number)
                continue
            if not page_text:
                continue
            chunks.append(
                SlideChunk(
                    slide_number=page_number,
                    text=page_text,
                    slide_title=f"Page {page_number}",
                    chunk_index=0,
                    source_type="page",
                )
            )
        return chunks


class SlideChunker:
    """Turns slide-level text into smaller semantic units."""

    def __init__(self, *, chunk_size: int = 500, chunk_overlap: int = 75) -> None:
        try:  # LangChain splitters moved to a standalone package in recent versions
            from langchain.text_splitter import RecursiveCharacterTextSplitter  # type: ignore
        except ModuleNotFoundError:  # pragma: no cover - executed in newer LangChain installs
            from langchain_text_splitters import RecursiveCharacterTextSplitter  # type: ignore

        self._splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            separators=["\n\n", "\n", ".", "!", "?", " "],
        )

    def chunk(self, slides: Sequence[SlideChunk]) -> List[SlideChunk]:
        """Split slide-level chunks into overlapping text segments."""
        processed: List[SlideChunk] = []
        for slide in slides:
            segments = self._splitter.split_text(slide.text)
            for index, segment in enumerate(segments):
                if not segment.strip():
                    continue
                processed.append(
                    SlideChunk(
                        slide_number=slide.slide_number,
                        text=segment.strip(),
                        slide_title=slide.slide_title,
                        chunk_index=index,
                        source_type=slide.source_type,
                    )
                )
        return processed


class EmbeddingService:
    """Generates embeddings using the configured provider."""

    def __init__(self, settings: Settings) -> None:
        if not settings.google_api_key:
            raise RuntimeError(
                "GOOGLE_API_KEY is required to generate embeddings. Update your .env with a valid key."
            )
        try:
            from langchain_google_genai import GoogleGenerativeAIEmbeddings  # type: ignore
        except ModuleNotFoundError as exc:  # pragma: no cover - import guard
            raise RuntimeError(
                "langchain-google-genai is required for embeddings. Install the dependency to continue."
            ) from exc

        self._settings = settings
        self._client = GoogleGenerativeAIEmbeddings(
            model=settings.google_embeddings_model_name,
            google_api_key=settings.google_api_key,
        )

    async def embed(self, texts: Sequence[str]) -> List[List[float]]:
        """Generate embeddings for a batch of texts using Google Generative AI embeddings."""
        payload = list(texts)
        if not payload:
            return []
        loop = asyncio.get_running_loop()
        return await loop.run_in_executor(None, self._client.embed_documents, payload)


class SlideIngestionPipeline:
    """Full orchestration for parsing, chunking, embedding, and indexing slides."""

    def __init__(
        self,
        settings: Settings,
        repository: Optional[PineconeRepository] = None,
        extractor: Optional[SlideExtractor] = None,
        pdf_extractor: Optional[PDFExtractor] = None,
        chunker: Optional[SlideChunker] = None,
        embedding_service: Optional[EmbeddingService] = None,
    ) -> None:
        self._settings = settings
        self._repository = repository or PineconeRepository(settings)
        self._pptx_extractor = extractor or SlideExtractor()
        self._pdf_extractor = pdf_extractor or PDFExtractor()
        self._chunker = chunker or SlideChunker()
        self._embedding_service = embedding_service or EmbeddingService(settings)

    async def ingest(
        self,
        *,
        document_id: str,
        file_bytes: bytes,
        filename: str | None = None,
        metadata: Optional[Dict[str, Any]] = None,
    ) -> IngestionResult:
        """Process a PPTX/PDF file through extraction, chunking, embedding, and Pinecone upsert."""
        extractor = self._select_extractor(filename)
        slides = extractor.extract(file_bytes)
        chunked = self._chunker.chunk(slides)

        if not chunked:
            logger.info("No text detected in presentation %s; skipping index", document_id)
            return IngestionResult(
                document_id=document_id,
                slide_count=len(slides),
                chunk_count=0,
                namespace=self._repository.namespace,
            )
        repo_dimension = getattr(self._repository, "dimension", None)
        dimension_validated = False
        batch_size = getattr(self._settings, "ingest_batch_size", 64) or 64
        base_metadata: Dict[str, Any] = dict(metadata or {})
        total_chunks = 0

        for batch in self._batched(chunked, batch_size):
            texts = [chunk.text for chunk in batch]
            # Generate embeddings and upsert to Pinecone so downstream chat/quiz can retrieve with citations.
            vectors = await self._embedding_service.embed(texts)
            if not vectors:
                continue
            if repo_dimension and not dimension_validated:
                embedding_dimension = len(vectors[0])
                if embedding_dimension != repo_dimension:
                    index_name = getattr(self._repository, "_index_name", "Pinecone index")
                    raise RuntimeError(
                        f"Embedding model produced dimension {embedding_dimension}, but Pinecone index "
                        f"{index_name} expects {repo_dimension}. "
                        "Ensure PINECONE_INDEX_DIMENSION matches the embedding model output and recreate/reconfigure "
                        "the Pinecone index if necessary."
                    )
                dimension_validated = True

            items: List[Dict[str, Any]] = []
            for chunk, embedding in zip(batch, vectors):
                payload = self._build_pinecone_payload(
                    chunk=chunk,
                    embedding=embedding,
                    base_metadata=base_metadata,
                    document_id=document_id,
                    snippet_chars=0,
                )
                items.append(payload)

            if items:
                self._repository.upsert(items)
                total_chunks += len(items)

        return IngestionResult(
            document_id=document_id,
            slide_count=len(slides),
            chunk_count=total_chunks,
            namespace=self._repository.namespace,
        )

    def _select_extractor(self, filename: str | None) -> SlideExtractor:
        """Choose the correct extractor based on filename; defaults to PPTX extractor."""
        if not filename:
            return self._pptx_extractor

        lowered = filename.lower()
        if lowered.endswith(".pptx"):
            return self._pptx_extractor
        if lowered.endswith(".pdf"):
            return self._pdf_extractor

        raise RuntimeError("Unsupported file type for ingestion; expected .pptx or .pdf")

    def delete_document(self, document_id: str) -> None:
        """Delete all vectors associated with a document id from Pinecone."""
        if not document_id:
            return
        self._repository.delete_document(document_id)

    @staticmethod
    def _batched(items: Sequence[SlideChunk], batch_size: int) -> Iterator[List[SlideChunk]]:
        """Yield fixed-size batches from a list of slide chunks."""
        if batch_size <= 0:
            batch_size = 1
        for start in range(0, len(items), batch_size):
            yield items[start : start + batch_size]

    def _build_pinecone_payload(
        self,
        *,
        chunk: SlideChunk,
        embedding: Sequence[float],
        base_metadata: Dict[str, Any],
        document_id: str,
        snippet_chars: int,
    ) -> Dict[str, Any]:
        """Build the Pinecone vector payload with merged metadata and optional snippet."""
        metadata_payload: Dict[str, Any] = {**base_metadata}
        metadata_payload.update(chunk.metadata())
        metadata_payload["document_id"] = document_id
        if snippet_chars > 0:
            snippet = chunk.text[:snippet_chars].strip()
            metadata_payload["text"] = snippet or chunk.text
        else:
            metadata_payload["text"] = chunk.text

        return {
            "id": f"{document_id}-s{chunk.slide_number}-c{chunk.chunk_index}",
            "values": list(embedding),
            "metadata": metadata_payload,
        }
