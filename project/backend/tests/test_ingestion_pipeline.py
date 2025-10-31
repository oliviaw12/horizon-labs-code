from __future__ import annotations

from io import BytesIO
from typing import Any, Dict, List, Sequence

import pytest
from pptx import Presentation

from clients.database.chat_repository import InMemoryChatRepository
from clients.ingestion.pipeline import (
    EmbeddingService,
    IngestionResult,
    SlideChunker,
    SlideExtractor,
    SlideIngestionPipeline,
)
from clients.llm.service import LLMService


def _build_presentation() -> bytes:
    ppt = Presentation()
    layout = ppt.slide_layouts[1]

    slide_one = ppt.slides.add_slide(layout)
    slide_one.shapes.title.text = "Introduction"
    slide_one.placeholders[1].text = "First bullet\nSecond bullet provides more context."

    slide_two = ppt.slides.add_slide(layout)
    slide_two.shapes.title.text = "Summary"
    slide_two.placeholders[1].text = "Takeaway one. Takeaway two continues the thought."

    buffer = BytesIO()
    ppt.save(buffer)
    return buffer.getvalue()


class FakeEmbeddingService(EmbeddingService):
    def __init__(self) -> None:
        pass

    async def embed(self, texts: Sequence[str]) -> List[List[float]]:
        return [[float(idx + 1)] for idx, _ in enumerate(texts)]


class CapturingRepository:
    def __init__(self) -> None:
        self.namespace = "slides"
        self.items: List[Dict[str, Any]] = []

    def upsert(self, items: Sequence[Dict[str, Any]]) -> None:
        self.items.extend(items)


@pytest.mark.asyncio
async def test_pipeline_extracts_chunks_and_writes_to_repository(test_settings) -> None:
    repository = CapturingRepository()
    pipeline = SlideIngestionPipeline(
        test_settings,
        repository=repository,
        extractor=SlideExtractor(),
        chunker=SlideChunker(chunk_size=60, chunk_overlap=0),
        embedding_service=FakeEmbeddingService(),
    )

    result = await pipeline.ingest(
        document_id="deck-1",
        file_bytes=_build_presentation(),
        metadata={"course": "math"},
    )

    assert isinstance(result, IngestionResult)
    assert result.document_id == "deck-1"
    assert result.slide_count == 2
    assert result.chunk_count == len(repository.items)
    assert repository.items, "expected upserted items"

    first = repository.items[0]
    assert first["metadata"]["slide_number"] == 1
    assert first["metadata"]["document_id"] == "deck-1"
    assert "text" in first["metadata"]


@pytest.mark.asyncio
async def test_llm_service_ingest_upload_uses_pipeline(monkeypatch, test_settings) -> None:
    test_settings.pinecone_api_key = "test"
    test_settings.pinecone_index_name = "index"

    service = LLMService(test_settings, repository=InMemoryChatRepository())

    class StubPipeline:
        def __init__(self) -> None:
            self.calls: List[Dict[str, Any]] = []

        async def ingest(self, **kwargs: Any) -> IngestionResult:
            self.calls.append(kwargs)
            return IngestionResult(
                document_id=kwargs["document_id"],
                slide_count=3,
                chunk_count=5,
                namespace="slides",
            )

    stub_pipeline = StubPipeline()
    service._ingestion_pipeline = stub_pipeline  # type: ignore[attr-defined]

    result = await service.ingest_upload(
        session_id="SESSION-123",
        file_bytes=b"binary",
        filename="Week1 Intro.pptx",
        metadata={"course": "math"},
    )

    assert result.chunk_count == 5
    assert stub_pipeline.calls, "ingestion pipeline should have been invoked"
    call = stub_pipeline.calls[0]
    assert call["metadata"]["session_id"] == "SESSION-123"
    assert call["metadata"]["source_filename"] == "Week1 Intro.pptx"
    assert call["document_id"].startswith("week1-intro")