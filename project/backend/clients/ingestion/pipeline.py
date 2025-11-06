from __future__ import annotations

import io
import logging
from dataclasses import dataclass
from typing import Any, Dict, List, Optional, Sequence
try:  # LangChain splitters moved to a standalone package in recent versions
    from langchain.text_splitter import RecursiveCharacterTextSplitter  # type: ignore
except ModuleNotFoundError:  # pragma: no cover - executed in newer LangChain installs
    from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_openai import OpenAIEmbeddings
from pptx import Presentation
from pypdf import PdfReader

from clients.database.pinecone import PineconeRepository
from clients.llm.settings import Settings

logger = logging.getLogger(__name__)


@dataclass
class SlideChunk:
    """A semantically meaningful chunk extracted from a document page or slide."""

    slide_number: int
    text: str
    slide_title: Optional[str]
    chunk_index: int
    source_type: str = "slide"

    def metadata(self) -> Dict[str, Any]:
        payload: Dict[str, Any] = {
            "slide_number": self.slide_number,
            "chunk_index": self.chunk_index,
            "slide_title": self.slide_title,
            "source_type": self.source_type,
        }
        if self.source_type == "page":
            payload["page_number"] = self.slide_number
        return payload


@dataclass
class IngestionResult:
    """Summary returned once the pipeline finishes ingesting a document."""

    document_id: str
    slide_count: int
    chunk_count: int
    namespace: str


class SlideExtractor:
    """Pulls raw text (and titles) from a slide deck."""

    def extract(self, file_bytes: bytes) -> List[SlideChunk]:
        presentation = Presentation(io.BytesIO(file_bytes))
        chunks: List[SlideChunk] = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            title = None
            if slide.shapes.title:
                title = slide.shapes.title.text.strip() or None

            text_fragments: List[str] = []
            for shape in slide.shapes:
                if not hasattr(shape, "text"):
                    continue
                payload = (shape.text or "").strip()
                if payload:
                    text_fragments.append(payload)

            slide_text = "\n".join(text_fragments).strip()
            if not slide_text:
                continue

            chunks.append(
                SlideChunk(
                    slide_number=slide_number,
                    text=slide_text,
                    slide_title=title,
                    chunk_index=0,
                    source_type="slide",
                )
            )
        return chunks


class PDFExtractor:
    """Extracts page-level text from a PDF document."""

    def extract(self, file_bytes: bytes) -> List[SlideChunk]:
        reader = PdfReader(io.BytesIO(file_bytes))
        chunks: List[SlideChunk] = []
        for page_number, page in enumerate(reader.pages, start=1):
            try:
                page_text = (page.extract_text() or "").strip()
            except Exception:  # pragma: no cover - defensive guard for uncommon PDFs
                logger.exception("Failed extracting text from PDF page %s", page_number)
                continue
            if not page_text:
                continue
            chunks.append(
                SlideChunk(
                    slide_number=page_number,
                    text=page_text,
                    slide_title=f"Page {page_number}",
                    chunk_index=0,
                    source_type="page",
                )
            )
        return chunks


class SlideChunker:
    """Turns slide-level text into smaller semantic units."""

    def __init__(self, *, chunk_size: int = 500, chunk_overlap: int = 75) -> None:
        self._splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            separators=["\n\n", "\n", ".", "!", "?", " "],
        )

    def chunk(self, slides: Sequence[SlideChunk]) -> List[SlideChunk]:
        processed: List[SlideChunk] = []
        for slide in slides:
            segments = self._splitter.split_text(slide.text)
            for index, segment in enumerate(segments):
                if not segment.strip():
                    continue
                processed.append(
                    SlideChunk(
                        slide_number=slide.slide_number,
                        text=segment.strip(),
                        slide_title=slide.slide_title,
                        chunk_index=index,
                        source_type=slide.source_type,
                    )
                )
        return processed


class EmbeddingService:
    """Generates embeddings using the configured provider."""

    def __init__(self, settings: Settings) -> None:
        self._settings = settings
        self._client = OpenAIEmbeddings(
            model=settings.embedding_model_name,
            openai_api_key=settings.openrouter_api_key,
            openai_api_base=settings.openrouter_base_url,
        )

    async def embed(self, texts: Sequence[str]) -> List[List[float]]:
        # LangChain OpenAIEmbeddings exposes async embedding via "aembed_documents".
        return await self._client.aembed_documents(list(texts))


class SlideIngestionPipeline:
    """Full orchestration for parsing, chunking, embedding, and indexing slides."""

    def __init__(
        self,
        settings: Settings,
        repository: Optional[PineconeRepository] = None,
        extractor: Optional[SlideExtractor] = None,
        pdf_extractor: Optional[PDFExtractor] = None,
        chunker: Optional[SlideChunker] = None,
        embedding_service: Optional[EmbeddingService] = None,
    ) -> None:
        self._settings = settings
        self._repository = repository or PineconeRepository(settings)
        self._pptx_extractor = extractor or SlideExtractor()
        self._pdf_extractor = pdf_extractor or PDFExtractor()
        self._chunker = chunker or SlideChunker()
        self._embedding_service = embedding_service or EmbeddingService(settings)

    async def ingest(
        self,
        *,
        document_id: str,
        file_bytes: bytes,
        filename: str | None = None,
        metadata: Optional[Dict[str, Any]] = None,
    ) -> IngestionResult:
        extractor = self._select_extractor(filename)
        slides = extractor.extract(file_bytes)
        chunked = self._chunker.chunk(slides)

        if not chunked:
            logger.info("No text detected in presentation %s; skipping index", document_id)
            return IngestionResult(
                document_id=document_id,
                slide_count=len(slides),
                chunk_count=0,
                namespace=self._repository.namespace,
            )

        texts = [chunk.text for chunk in chunked]
        vectors = await self._embedding_service.embed(texts)

        items: List[Dict[str, Any]] = []
        for chunk, embedding in zip(chunked, vectors):
            payload = {
                "id": f"{document_id}-s{chunk.slide_number}-c{chunk.chunk_index}",
                "values": embedding,
                "metadata": {
                    **(metadata or {}),
                    **chunk.metadata(),
                    "document_id": document_id,
                    "text": chunk.text,
                },
            }
            items.append(payload)

        self._repository.upsert(items)

        return IngestionResult(
            document_id=document_id,
            slide_count=len(slides),
            chunk_count=len(items),
            namespace=self._repository.namespace,
        )

    def _select_extractor(self, filename: str | None) -> SlideExtractor:
        if not filename:
            return self._pptx_extractor

        lowered = filename.lower()
        if lowered.endswith(".pptx"):
            return self._pptx_extractor
        if lowered.endswith(".pdf"):
            return self._pdf_extractor

        raise RuntimeError("Unsupported file type for ingestion; expected .pptx or .pdf")